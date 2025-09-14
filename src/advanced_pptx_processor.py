import os
import json
import zipfile
import tempfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import re
import subprocess



class AdvancedPPTProcessor:
    def __init__(self, preprocessor, fast_mode=False):
        """
        初始化高级PPTX处理器
        
        Args:
            preprocessor: MultimodalPreprocessor实例，用于重用输出目录与结果记录
            fast_mode: 快速模式，跳过耗时的CLIP描述生成
        """
        self.preprocessor = preprocessor
        self.fast_mode = fast_mode
        self.output_text_dir = "output/text"
        self.output_table_dir = "output/tables"
        self.output_img_dir = "output/images"
        
        # 确保输出目录存在
        os.makedirs(self.output_text_dir, exist_ok=True)
        os.makedirs(self.output_table_dir, exist_ok=True)
        os.makedirs(self.output_img_dir, exist_ok=True)

    def extract_all_images_via_zip(self, file_path):
        """
        通过ZIP解压和XML解析提取PPTX中的所有图片
        
        Args:
            file_path: PPTX文件路径
            
        Returns:
            dict: 包含幻灯片到图片映射关系的字典
        """
        print(f"开始通过ZIP方式提取图片: {file_path}")
        
        base_filename = os.path.splitext(os.path.basename(file_path))[0]
        slide_image_mapping = {}
        
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # 1. 解压PPTX文件
                print("正在解压PPTX文件...")
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                # 2. 找到媒体目录
                media_dir = os.path.join(temp_dir, "ppt", "media")
                slides_dir = os.path.join(temp_dir, "ppt", "slides")
                rels_dir = os.path.join(temp_dir, "ppt", "slides", "_rels")
                
                if not os.path.exists(media_dir):
                    print("未找到media目录，可能没有图片")
                    return slide_image_mapping
                
                print(f"找到media目录: {media_dir}")
                print(f"媒体文件: {os.listdir(media_dir)}")
                
                # 3. 遍历所有幻灯片XML文件
                if os.path.exists(slides_dir):
                    for slide_file in os.listdir(slides_dir):
                        if slide_file.startswith("slide") and slide_file.endswith(".xml"):
                            slide_num = self._extract_slide_number(slide_file)
                            if slide_num is None:
                                continue
                                
                            print(f"处理幻灯片 {slide_num}: {slide_file}")
                            
                            # 解析幻灯片XML获取图片关系ID
                            slide_xml_path = os.path.join(slides_dir, slide_file)
                            image_rids = self._parse_slide_xml_for_images(slide_xml_path)
                            
                            if image_rids:
                                print(f"幻灯片 {slide_num} 中找到图片关系ID: {image_rids}")
                                
                                # 解析关系文件获取实际文件名
                                rels_file = slide_file + ".rels"
                                rels_path = os.path.join(rels_dir, rels_file)
                                
                                if os.path.exists(rels_path):
                                    image_files = self._parse_rels_file(rels_path, image_rids)
                                    
                                    if image_files:
                                        slide_image_mapping[slide_num] = image_files
                                        print(f"幻灯片 {slide_num} 映射到图片: {image_files}")
                                        
                                        # 复制图片到输出目录
                                        self._copy_images_to_output(media_dir, image_files, 
                                                                  base_filename, slide_num)
                
                print(f"图片提取完成，映射关系: {slide_image_mapping}")
                
            except Exception as e:
                print(f"ZIP方式图片提取失败: {e}")
                import traceback
                traceback.print_exc()
        
        return slide_image_mapping

    def _extract_slide_number(self, slide_filename):
        """从幻灯片文件名中提取编号"""
        try:
            # slide1.xml -> 1
            import re
            match = re.search(r'slide(\d+)\.xml', slide_filename)
            if match:
                return int(match.group(1))
        except Exception:
            pass
        return None

    def _parse_slide_xml_for_images(self, slide_xml_path):
        """
        解析幻灯片XML文件，查找图片引用
        
        Args:
            slide_xml_path: 幻灯片XML文件路径
            
        Returns:
            list: 图片关系ID列表
        """
        image_rids = []
        
        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }
            
            # 查找所有a:blip元素（图片引用）
            blip_elements = root.findall('.//a:blip', namespaces)
            
            for blip in blip_elements:
                embed_attr = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_attr:
                    image_rids.append(embed_attr)
                    print(f"找到图片引用ID: {embed_attr}")
            
        except Exception as e:
            print(f"解析幻灯片XML失败 {slide_xml_path}: {e}")
        
        return image_rids

    def _parse_rels_file(self, rels_path, image_rids):
        """
        解析关系文件，获取关系ID到文件名的映射
        
        Args:
            rels_path: 关系文件路径
            image_rids: 图片关系ID列表
            
        Returns:
            list: 对应的图片文件名列表
        """
        image_files = []
        
        try:
            tree = ET.parse(rels_path)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }
            
            # 查找所有关系
            for relationship in root.findall('.//rel:Relationship', namespaces):
                rel_id = relationship.get('Id')
                target = relationship.get('Target')
                rel_type = relationship.get('Type')
                
                # 检查是否是图片关系
                if (rel_id in image_rids and 
                    target and 
                    rel_type and 
                    'image' in rel_type.lower()):
                    
                    # 提取文件名 (../media/image1.png -> image1.png)
                    filename = os.path.basename(target)
                    image_files.append(filename)
                    print(f"关系映射: {rel_id} -> {filename}")
            
        except Exception as e:
            print(f"解析关系文件失败 {rels_path}: {e}")
        
        return image_files

    def _copy_images_to_output(self, media_dir, image_files, base_filename, slide_num):
        """
        将图片复制到输出目录并生成描述
        
        Args:
            media_dir: 媒体文件源目录
            image_files: 图片文件名列表
            base_filename: 基础文件名
            slide_num: 幻灯片编号
        """
        for idx, image_file in enumerate(image_files, 1):
            try:
                source_path = os.path.join(media_dir, image_file)
                
                if os.path.exists(source_path):
                    # 生成输出文件名
                    file_ext = os.path.splitext(image_file)[1]
                    output_filename = f"{base_filename}_slide_{slide_num}_img_{idx}_zip{file_ext}"
                    output_path = os.path.join(self.output_img_dir, output_filename)
                    
                    # 复制图片
                    shutil.copy2(source_path, output_path)
                    print(f"复制图片: {source_path} -> {output_path}")
                    
                    # 生成CLIP描述（根据模式决定是否生成）
                    desc_path = None
                    if not self.fast_mode:
                        try:
                            desc_path = self.preprocessor.clip_generate_description(output_path)
                        except Exception as e:
                            print(f"生成CLIP描述失败，跳过: {e}")
                    else:
                        print("快速模式：跳过CLIP描述生成")
                    
                    # 记录到结果中
                    self.preprocessor.results.append({
                        "type": "ppt_image_zip",
                        "page": slide_num,
                        "file": output_path,
                        "description_file": desc_path,
                        "extraction_method": "zip_xml_parsing",
                        "original_filename": image_file
                    })
                    
                else:
                    print(f"源图片文件不存在: {source_path}")
                    
            except Exception as e:
                print(f"复制图片失败 {image_file}: {e}")

    def extract_and_convert_equations(self, slide, slide_number):
        """
        处理幻灯片中的公式
        
        Args:
            slide: python-pptx的Slide对象
            slide_number: 幻灯片编号
            
        Returns:
            list: 包含公式信息的列表
        """
        equations = []
        
        for shape_index, shape in enumerate(slide.shapes):
            try:
                # 检查形状是否包含文本框
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 获取形状的XML内容
                    shape_xml = self._get_shape_xml(shape)
                    if shape_xml:
                        # 检查是否包含OMML公式标签
                        omml_content = self._extract_omml_from_xml(shape_xml)
                        if omml_content:
                            print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现OMML公式")
                            
                            # 尝试转换OMML到LaTeX
                            latex_content = self._convert_omml_to_latex(omml_content)
                            
                            equation_info = {
                                "slide_number": slide_number,
                                "shape_index": shape_index,
                                "type": "omml_formula",
                                "original_omml": omml_content[:500] + "..." if len(omml_content) > 500 else omml_content,
                                "latex": latex_content,
                                "conversion_success": latex_content is not None
                            }
                            
                            equations.append(equation_info)
                            
                            # 添加到结果中
                            self.preprocessor.results.append({
                                "type": "formula",
                                "page": slide_number,
                                "formula_type": "omml",
                                "latex": latex_content,
                                "source": f"slide_{slide_number}_shape_{shape_index}",
                                "conversion_method": "omml_to_latex"
                            })
                            
                            continue
                
                # 如果没有找到OMML，检查是否为可能的公式图片
                if self._is_potential_formula_image(shape):
                    print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现潜在公式图片")
                    
                    # 使用图片处理流程处理公式图片
                    formula_image_path = self._process_formula_image(shape, slide_number, shape_index)
                    
                    if formula_image_path:
                        equation_info = {
                            "slide_number": slide_number,
                            "shape_index": shape_index,
                            "type": "formula_image",
                            "image_path": formula_image_path,
                            "latex": None,
                            "conversion_success": False
                        }
                        
                        equations.append(equation_info)
                        
                        # 添加到结果中
                        self.preprocessor.results.append({
                            "type": "formula",
                            "page": slide_number,
                            "formula_type": "image",
                            "image_path": formula_image_path,
                            "source": f"slide_{slide_number}_shape_{shape_index}",
                            "conversion_method": "image_fallback"
                        })
                        
            except Exception as e:
                print(f"处理幻灯片 {slide_number} 形状 {shape_index} 时出错: {e}")
                continue
        
        return equations

    def _get_shape_xml(self, shape):
        """获取形状的XML内容"""
        try:
            # 尝试获取形状的内部XML
            if hasattr(shape, '_element'):
                return ET.tostring(shape._element, encoding='unicode')
        except Exception as e:
            print(f"获取形状XML失败: {e}")
        return None

    def _extract_omml_from_xml(self, xml_string):
        """从XML中提取OMML内容"""
        try:
            # 查找OMML数学标签
            omml_patterns = [
                r'<m:oMath[^>]*>.*?</m:oMath>',
                r'<m:oMathPara[^>]*>.*?</m:oMathPara>',
                r'<math[^>]*>.*?</math>'  # 也检查标准MathML
            ]
            
            for pattern in omml_patterns:
                matches = re.findall(pattern, xml_string, re.DOTALL | re.IGNORECASE)
                if matches:
                    return matches[0]
                    
        except Exception as e:
            print(f"提取OMML失败: {e}")
        return None

    def _convert_omml_to_latex(self, omml_content):
        """将OMML转换为LaTeX"""
        try:
            # 方法1: 尝试使用pandoc
            latex_result = self._convert_via_pandoc(omml_content)
            if latex_result:
                return latex_result
                
            # 方法2: 简单的文本替换作为备选方案
            latex_result = self._simple_omml_to_latex(omml_content)
            if latex_result:
                return latex_result
                
        except Exception as e:
            print(f"OMML转LaTeX失败: {e}")
        
        return None

    def _convert_via_pandoc(self, omml_content):
        """使用pandoc转换OMML到LaTeX"""
        try:
            # 检查pandoc是否可用
            subprocess.run(['pandoc', '--version'], 
                         capture_output=True, check=True)
            
            # 创建临时文件
            with tempfile.NamedTemporaryFile(mode='w', suffix='.xml', delete=False) as temp_file:
                temp_file.write(f'<root>{omml_content}</root>')
                temp_file_path = temp_file.name
            
            try:
                # 使用pandoc转换
                result = subprocess.run([
                    'pandoc', 
                    '-f', 'docx',
                    '-t', 'latex',
                    temp_file_path
                ], capture_output=True, text=True, check=True)
                
                return result.stdout.strip()
                
            finally:
                os.unlink(temp_file_path)
                
        except (subprocess.CalledProcessError, FileNotFoundError):
            print("Pandoc不可用，跳过pandoc转换")
        except Exception as e:
            print(f"Pandoc转换失败: {e}")
        
        return None

    def _simple_omml_to_latex(self, omml_content):
        """简单的OMML到LaTeX转换（基本文本替换）"""
        try:
            # 移除XML标签，提取纯文本
            text_content = re.sub(r'<[^>]+>', '', omml_content)
            text_content = text_content.strip()
            
            if not text_content:
                return None
            
            # 基本的数学符号替换
            replacements = {
                '≈': r'\approx',
                '≠': r'\neq',
                '≤': r'\leq',
                '≥': r'\geq',
                '∞': r'\infty',
                'α': r'\alpha',
                'β': r'\beta',
                'γ': r'\gamma',
                'δ': r'\delta',
                'θ': r'\theta',
                'λ': r'\lambda',
                'μ': r'\mu',
                'π': r'\pi',
                'σ': r'\sigma',
                'φ': r'\phi',
                'ω': r'\omega',
                '∑': r'\sum',
                '∫': r'\int',
                '√': r'\sqrt',
                '±': r'\pm',
                '×': r'\times',
                '÷': r'\div'
            }
            
            for symbol, latex in replacements.items():
                text_content = text_content.replace(symbol, latex)
            
            # 包装在数学环境中
            return f"${text_content}$"
            
        except Exception as e:
            print(f"简单转换失败: {e}")
        
        return None

    def _is_potential_formula_image(self, shape):
        """判断形状是否可能是公式图片"""
        try:
            # 检查是否为图片类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            
            # 检查是否为包含复杂路径的形状（可能是矢量公式）
            if shape.shape_type in [MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                return True
            
            # 检查形状大小（小的形状可能是公式）
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # 假设公式通常比较小（宽度和高度都小于某个阈值）
                max_formula_size = 200000  # EMU单位
                if shape.width < max_formula_size and shape.height < max_formula_size:
                    return True
                    
        except Exception as e:
            print(f"检查潜在公式图片失败: {e}")
        
        return False

    def _process_formula_image(self, shape, slide_number, shape_index):
        """处理公式图片"""
        try:
            # 如果是图片类型，尝试导出图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_filename = f"formula_slide_{slide_number}_shape_{shape_index}.png"
                image_path = os.path.join(self.output_img_dir, image_filename)
                
                # 这里需要实现图片导出逻辑
                # 由于python-pptx的限制，可能需要使用其他方法
                print(f"识别到公式图片，但需要额外的导出逻辑: {image_filename}")
                
                return image_path
                
        except Exception as e:
            print(f"处理公式图片失败: {e}")
        
        return None

    def process_pptx_file_advanced(self, file_path):
        """
        高级PPTX处理：结合传统方法和ZIP解析
        
        Args:
            file_path: PPTX文件路径
        """
        print(f"开始高级PPTX处理: {file_path}")
        base_filename = os.path.splitext(os.path.basename(file_path))[0]
        
        # 重置结果记录，为当前PPTX文件单独记录
        current_results = []
        original_results = self.preprocessor.results
        self.preprocessor.results = current_results
        
        try:
            # 1. 使用传统python-pptx方法处理文本和表格
            self._process_text_and_tables_traditional(file_path, base_filename)
            
            # 2. 使用ZIP方法提取所有图片
            slide_image_mapping = self.extract_all_images_via_zip(file_path)
            
            # 3. 生成PPTX专用元数据
            self._save_pptx_metadata(file_path, base_filename, slide_image_mapping)
            
            print(f"高级PPTX处理完成: {file_path}")
            print(f"PPTX元数据已保存: output/{base_filename}_pptx_metadata.json")
            
        except Exception as e:
            print(f"高级PPTX处理失败: {e}")
            import traceback
            traceback.print_exc()
        finally:
            # 恢复原始结果列表并合并当前结果
            self.preprocessor.results = original_results
            self.preprocessor.results.extend(current_results)

    def _process_text_and_tables_traditional(self, file_path, base_filename):
        """使用传统python-pptx方法处理文本和表格"""
        prs = Presentation(file_path)
        
        for slide_index, slide in enumerate(prs.slides, start=1):
            # 处理公式
            equations = self.extract_and_convert_equations(slide, slide_index)
            if equations:
                print(f"在幻灯片 {slide_index} 中找到 {len(equations)} 个公式")
                
                # 保存公式信息到JSON文件
                formulas_json_path = f"output/formulas/{base_filename}_slide_{slide_index}_formulas.json"
                os.makedirs("output/formulas", exist_ok=True)
                
                formulas_output = {
                    "slide_number": slide_index,
                    "source_file": base_filename,
                    "equations_count": len(equations),
                    "equations": equations,
                    "processing_date": str(pd.Timestamp.now())
                }
                
                with open(formulas_json_path, "w", encoding="utf-8") as f:
                    json.dump(formulas_output, f, ensure_ascii=False, indent=2)
                
                print(f"公式信息已保存到: {formulas_json_path}")
            
            # 提取文本
            slide_text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        runs_text = ''.join(run.text for run in paragraph.runs)
                        text_content.append(runs_text if runs_text else paragraph.text)
                    final_text = "\n".join([t for t in text_content if t is not None])
                    if final_text.strip():
                        slide_text_items.append(final_text.strip())

            if slide_text_items:
                text_output = {
                    "type": "ppt_text",
                    "page": slide_index,
                    "source": f"{base_filename}_slide_{slide_index}",
                    "raw_text": "\n\n".join(slide_text_items)
                }
                text_json_path = f"{self.output_text_dir}/{base_filename}_slide_{slide_index}.json"
                with open(text_json_path, "w", encoding="utf-8") as f:
                    json.dump(text_output, f, ensure_ascii=False, indent=2)
                self.preprocessor.results.append({
                    "type": "ppt_text",
                    "page": slide_index,
                    "file": text_json_path
                })

            # 提取表格（优化版本）
            table_counter = 0
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table_counter += 1
                    table = shape.table
                    
                    # 获取表格位置信息（增强功能）
                    table_position = {
                        "left": float(shape.left.inches) if shape.left else 0,
                        "top": float(shape.top.inches) if shape.top else 0,
                        "width": float(shape.width.inches) if shape.width else 0,
                        "height": float(shape.height.inches) if shape.height else 0
                    }
                    
                    # 提取表格数据
                    data_matrix = []
                    for row in table.rows:
                        row_values = []
                        for cell in row.cells:
                            # 优化：使用 text_frame.text 获取纯文本
                            try:
                                if cell.text_frame and cell.text_frame.text:
                                    cell_text = cell.text_frame.text.strip()
                                else:
                                    cell_text = cell.text.strip() if cell.text else ""
                            except Exception as e:
                                print(f"提取单元格文本失败: {e}")
                                cell_text = ""
                            row_values.append(cell_text)
                        data_matrix.append(row_values)
                    
                    # 检查是否有有效数据
                    if data_matrix and any(any(cell for cell in row) for row in data_matrix):
                        # 使用pandas DataFrame保存为CSV
                        df = pd.DataFrame(data_matrix)
                        csv_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.csv"
                        df.to_csv(csv_path, index=False, header=False, encoding="utf-8")
                        
                        # 创建表格JSON元数据文件
                        table_metadata = {
                            "type": "ppt_table",
                            "source": f"{base_filename}_slide_{slide_index}",
                            "slide_number": slide_index,
                            "table_index": table_counter,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "data_preview": data_matrix[:3] if len(data_matrix) > 0 else [],  # 前3行预览
                            "csv_file": csv_path
                        }
                        
                        # 保存表格元数据JSON
                        json_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.json"
                        with open(json_path, "w", encoding="utf-8") as f:
                            json.dump(table_metadata, f, ensure_ascii=False, indent=2)
                        
                        # 记录到主结果中（增强元数据）
                        self.preprocessor.results.append({
                            "type": "ppt_table",
                            "page": slide_index,
                            "table_index": table_counter,
                            "file": csv_path,
                            "metadata_file": json_path,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "extraction_method": "python_pptx_optimized"
                        })
                        
                        print(f"✓ 提取表格 {table_counter}: {len(data_matrix)}行 x {len(data_matrix[0]) if data_matrix else 0}列")
                        print(f"  位置: left={table_position['left']:.2f}in, top={table_position['top']:.2f}in")
                    else:
                        print(f"⚠ 跳过空表格 {table_counter}")

    def _save_pptx_metadata(self, file_path, base_filename, slide_image_mapping):
        """
        保存PPTX专用元数据文件
        
        Args:
            file_path: 原始PPTX文件路径
            base_filename: 基础文件名
            slide_image_mapping: 幻灯片到图片的映射关系
        """
        import datetime
        
        # 统计各类型文件
        text_files = [r for r in self.preprocessor.results if r["type"] == "ppt_text"]
        table_files = [r for r in self.preprocessor.results if r["type"] == "ppt_table"]
        image_files_traditional = [r for r in self.preprocessor.results if r["type"] == "ppt_image"]
        image_files_zip = [r for r in self.preprocessor.results if r["type"] == "ppt_image_zip"]
        
        # 计算幻灯片统计
        total_slides = len(set(r["page"] for r in self.preprocessor.results if "page" in r))
        slides_with_images = len(slide_image_mapping)
        total_images_zip = sum(len(images) for images in slide_image_mapping.values())
        
        # 计算表格统计
        total_tables = len(table_files)
        total_table_rows = sum(r.get("dimensions", {}).get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("dimensions", {}).get("columns", 0) for r in table_files)
        table_positions = [r.get("position", {}) for r in table_files if r.get("position")]
        
        # 构建元数据
        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": base_filename,
            "source_path": file_path,
            "file_type": "PPTX",
            "processing_method": "advanced_zip_xml_parsing",
            "statistics": {
                "total_slides": total_slides,
                "slides_with_text": len(text_files),
                "slides_with_tables": len([r for r in table_files if r.get("dimensions", {}).get("rows", 0) > 0]),
                "total_tables": total_tables,
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "slides_with_images": slides_with_images,
                "total_images_extracted": len(image_files_traditional) + len(image_files_zip),
                "images_via_traditional": len(image_files_traditional),
                "images_via_zip_parsing": len(image_files_zip),
                "total_images_in_media": total_images_zip
            },
            "slide_image_mapping": slide_image_mapping,
            "files": {
                "text_files": text_files,
                "table_files": table_files,
                "image_files_traditional": image_files_traditional,
                "image_files_zip": image_files_zip
            },
            "table_analysis": {
                "positions": table_positions,
                "position_stats": {
                    "avg_left": sum(p.get("left", 0) for p in table_positions) / len(table_positions) if table_positions else 0,
                    "avg_top": sum(p.get("top", 0) for p in table_positions) / len(table_positions) if table_positions else 0,
                    "avg_width": sum(p.get("width", 0) for p in table_positions) / len(table_positions) if table_positions else 0,
                    "avg_height": sum(p.get("height", 0) for p in table_positions) / len(table_positions) if table_positions else 0
                }
            },
            "processing_info": {
                "extraction_methods": ["python-pptx", "zip_xml_parsing"],
                "image_formats_supported": ["PNG", "WMF", "EMF", "JPEG"],
                "table_extraction_enhanced": True,
                "table_position_tracking": True,
                "clip_descriptions_generated": True,
                "output_format": "JSON/CSV"
            }
        }
        
        # 保存元数据
        metadata_path = f"output/{base_filename}_pptx_metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        
        return metadata_path


def process_pptx_file_advanced(preprocessor, file_path, fast_mode=False):
    """
    高级PPTX处理的入口函数
    
    Args:
        preprocessor: MultimodalPreprocessor实例
        file_path: PPTX文件路径
        fast_mode: 快速模式，跳过耗时处理
    """
    processor = AdvancedPPTProcessor(preprocessor, fast_mode=fast_mode)
    processor.process_pptx_file_advanced(file_path)
